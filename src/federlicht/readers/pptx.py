from __future__ import annotations

import io
from pathlib import Path
from typing import Optional

from ..utils.strings import slugify_url

try:
    import pptx  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
except Exception:
    pptx = None
    MSO_SHAPE_TYPE = None


def _pptx_available() -> bool:
    return pptx is not None


def _pillow_image() -> Optional[object]:
    try:
        from PIL import Image  # type: ignore
    except Exception:
        return None
    return Image


def _emu_to_px(value: int) -> int:
    # 1 inch = 914400 EMU, assume 96 DPI -> 1 px = 9525 EMU.
    return max(int(value / 9525), 0)


def _slide_title(slide) -> str:
    try:
        title = slide.shapes.title
    except Exception:
        title = None
    if title is None:
        return ""
    text = getattr(title, "text", "") or ""
    return text.strip()


def _extract_table_text(shape) -> list[str]:
    try:
        table = shape.table
    except Exception:
        return []
    rows: list[str] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            value = (cell.text or "").strip()
            cells.append(value)
        if any(cells):
            rows.append(" | ".join(cells))
    return rows


def _extract_chart_text(shape, max_points: int = 6) -> list[str]:
    try:
        chart = shape.chart
    except Exception:
        return []
    lines: list[str] = []
    try:
        title = chart.chart_title.text_frame.text if chart.has_title else ""
    except Exception:
        title = ""
    if title:
        lines.append(f"Chart title: {title.strip()}")
    try:
        series = list(chart.series)
    except Exception:
        series = []
    for idx, ser in enumerate(series, start=1):
        name = (getattr(ser, "name", None) or f"Series {idx}").strip()
        values = []
        try:
            values = list(ser.values)
        except Exception:
            values = []
        if values:
            preview = ", ".join(str(v) for v in values[:max_points])
            if len(values) > max_points:
                preview = f"{preview}, …"
            lines.append(f"{name}: {preview}")
        else:
            lines.append(f"{name}")
    return lines


def _extract_text_frame(shape) -> list[str]:
    try:
        frame = shape.text_frame
    except Exception:
        return []
    lines: list[str] = []
    for para in frame.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        indent = "  " * int(getattr(para, "level", 0) or 0)
        prefix = "- " if indent or len(lines) else ""
        lines.append(f"{indent}{prefix}{text}")
    return lines


def _source_path_under_run(path: Path, run_dir: Path) -> str:
    try:
        rel = path.resolve().relative_to(run_dir.resolve()).as_posix()
        return f"./{rel}"
    except Exception:
        return path.resolve().as_posix()


def _shape_type_name(shape) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if MSO_SHAPE_TYPE is not None and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if getattr(shape, "has_text_frame", False):
        return "text"
    return "unknown"


def _shape_text_payload(shape, shape_type: str) -> str:
    lines: list[str] = []
    if shape_type == "table":
        lines = _extract_table_text(shape)
    elif shape_type == "chart":
        lines = _extract_chart_text(shape)
    elif shape_type == "text":
        lines = _extract_text_frame(shape)
    return "\n".join(line for line in lines if str(line).strip()).strip()


def extract_pptx_slide_contract(
    pptx_path: Path,
    run_dir: Path,
    *,
    max_slides: int = 0,
    start_slide: int = 0,
    include_notes: bool = True,
    max_text_chars_per_shape: int = 500,
) -> dict:
    source_path = _source_path_under_run(pptx_path, run_dir)
    if not _pptx_available():
        return {
            "schema_version": "pptx_ingest.v1",
            "reader": "python-pptx",
            "available": False,
            "error": "python-pptx is not installed",
            "source_path": source_path,
            "slide_count_total": 0,
            "extracted_slide_count": 0,
            "truncated": False,
            "slides": [],
        }

    prs = pptx.Presentation(pptx_path)  # type: ignore[attr-defined]
    total = len(prs.slides)
    if total <= 0:
        return {
            "schema_version": "pptx_ingest.v1",
            "reader": "python-pptx",
            "available": True,
            "source_path": source_path,
            "slide_count_total": 0,
            "extracted_slide_count": 0,
            "slide_start": 0,
            "slide_end": 0,
            "truncated": False,
            "slides": [],
        }

    start_slide = max(0, min(start_slide, max(0, total - 1)))
    if max_slides <= 0:
        count = max(0, total - start_slide)
    else:
        count = min(max_slides, total - start_slide)

    slides: list[dict] = []
    for idx in range(start_slide, start_slide + count):
        slide = prs.slides[idx]
        slide_no = idx + 1
        slide_id = f"slide-{slide_no:03d}"
        slide_anchor = f"{source_path}#slide-{slide_no}"
        elements: list[dict] = []
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            shape_type = _shape_type_name(shape)
            text_payload = _shape_text_payload(shape, shape_type)
            if max_text_chars_per_shape > 0 and text_payload:
                text_payload = text_payload[:max_text_chars_per_shape]
            content_kind = "image" if shape_type == "picture" else "text"
            element: dict[str, object] = {
                "element_id": f"{slide_id}-shape-{shape_idx:03d}",
                "shape_index": shape_idx,
                "shape_type": shape_type,
                "content_kind": content_kind,
                "source_path": source_path,
                "anchor": f"{slide_anchor}-shape-{shape_idx}",
            }
            if text_payload:
                element["text"] = text_payload
            if shape_type == "picture":
                width = _emu_to_px(int(getattr(shape, "width", 0) or 0))
                height = _emu_to_px(int(getattr(shape, "height", 0) or 0))
                element["image"] = {"width_px": max(0, width), "height_px": max(0, height)}
            if shape_type == "unknown" and not text_payload:
                continue
            elements.append(element)

        if include_notes:
            try:
                notes = slide.notes_slide.notes_text_frame.text  # type: ignore[attr-defined]
            except Exception:
                notes = ""
            note_text = str(notes or "").strip()
            if max_text_chars_per_shape > 0 and note_text:
                note_text = note_text[:max_text_chars_per_shape]
            if note_text:
                elements.append(
                    {
                        "element_id": f"{slide_id}-notes",
                        "shape_index": len(elements) + 1,
                        "shape_type": "notes",
                        "content_kind": "text",
                        "source_path": source_path,
                        "anchor": f"{slide_anchor}-notes",
                        "text": note_text,
                    }
                )

        slides.append(
            {
                "slide_id": slide_id,
                "slide_index": slide_no,
                "slide_title": _slide_title(slide),
                "source_path": source_path,
                "anchor": slide_anchor,
                "elements": elements,
            }
        )

    return {
        "schema_version": "pptx_ingest.v1",
        "reader": "python-pptx",
        "available": True,
        "source_path": source_path,
        "slide_count_total": total,
        "extracted_slide_count": len(slides),
        "slide_start": start_slide + 1,
        "slide_end": start_slide + count,
        "truncated": (start_slide + count) < total,
        "slides": slides,
    }


def read_pptx_text(
    pptx_path: Path,
    max_slides: int,
    max_chars: int,
    start_slide: int = 0,
    include_notes: bool = True,
) -> str:
    if not _pptx_available():
        return "python-pptx is not installed. Cannot read PPTX."
    prs = pptx.Presentation(pptx_path)  # type: ignore[attr-defined]
    total = len(prs.slides)
    start_slide = max(0, min(start_slide, max(0, total - 1)))
    if max_slides <= 0:
        count = max(0, total - start_slide)
    else:
        count = min(max_slides, total - start_slide)

    chunks: list[str] = []
    for idx in range(start_slide, start_slide + count):
        slide = prs.slides[idx]
        slide_no = idx + 1
        title = _slide_title(slide)
        header = f"Slide {slide_no}"
        if title:
            header = f"{header}: {title}"
        lines = [header]
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table_lines = _extract_table_text(shape)
                if table_lines:
                    lines.append("Table:")
                    lines.extend(f"  {row}" for row in table_lines)
                continue
            if getattr(shape, "has_chart", False):
                chart_lines = _extract_chart_text(shape)
                if chart_lines:
                    lines.append("Chart:")
                    lines.extend(f"  {row}" for row in chart_lines)
                continue
            if getattr(shape, "has_text_frame", False):
                text_lines = _extract_text_frame(shape)
                if text_lines:
                    lines.extend(text_lines)
        if include_notes:
            try:
                notes = slide.notes_slide.notes_text_frame.text  # type: ignore[attr-defined]
            except Exception:
                notes = ""
            notes = (notes or "").strip()
            if notes:
                lines.append("Notes:")
                lines.append(notes)
        if len(lines) > 1:
            chunks.append("\n".join(lines))

    text = "\n\n".join(chunks)
    note = ""
    if start_slide + count < total:
        first = start_slide + 1
        last = start_slide + count
        note = (
            f"\n\n[note] PPTX scan truncated: slides {first}-{last} of {total}. "
            "Increase --max-pptx-slides or use start_slide to read more."
        )
    if max_chars > 0:
        if note:
            remaining = max_chars - len(note)
            if remaining <= 0:
                return note[:max_chars]
            return f"{text[:remaining]}{note}"
        return text[:max_chars]
    return f"{text}{note}"


def extract_pptx_images(
    pptx_path: Path,
    output_dir: Path,
    run_dir: Path,
    max_per_pptx: int,
    min_area: int,
) -> list[dict]:
    if not _pptx_available():
        return []
    output_dir.mkdir(parents=True, exist_ok=True)
    prs = pptx.Presentation(pptx_path)  # type: ignore[attr-defined]
    records: list[dict] = []
    candidates: list[dict] = []
    pptx_rel = f"./{pptx_path.relative_to(run_dir).as_posix()}"
    Image = _pillow_image()
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if MSO_SHAPE_TYPE is None or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
                blob = image.blob
                ext = (image.ext or "png").lower()
            except Exception:
                continue
            width = 0
            height = 0
            if Image:
                try:
                    with Image.open(io.BytesIO(blob)) as pil:
                        width, height = pil.size
                except Exception:
                    width = _emu_to_px(int(getattr(shape, "width", 0) or 0))
                    height = _emu_to_px(int(getattr(shape, "height", 0) or 0))
            else:
                width = _emu_to_px(int(getattr(shape, "width", 0) or 0))
                height = _emu_to_px(int(getattr(shape, "height", 0) or 0))
            if width and height and width * height < min_area:
                continue
            tag = f"{pptx_rel}#s{slide_idx}-{len(candidates) + 1}"
            candidates.append(
                {
                    "pptx_path": pptx_rel,
                    "slide": slide_idx,
                    "width": width,
                    "height": height,
                    "area": width * height,
                    "tag": tag,
                    "ext": ext,
                    "image": blob,
                    "slide_title": _slide_title(slide),
                }
            )

    if candidates:
        candidates.sort(key=lambda item: item["area"], reverse=True)
        for candidate in candidates[:max_per_pptx]:
            name = f"{slugify_url(candidate['tag'])}.{candidate['ext']}"
            img_path = output_dir / name
            if not img_path.exists():
                try:
                    with img_path.open("wb") as handle:
                        handle.write(candidate["image"])
                except Exception:
                    continue
            img_rel = f"./{img_path.relative_to(run_dir).as_posix()}"
            records.append(
                {
                    "pptx_path": candidate["pptx_path"],
                    "image_path": img_rel,
                    "slide": candidate["slide"],
                    "width": candidate["width"],
                    "height": candidate["height"],
                    "method": "embedded",
                    "slide_title": candidate.get("slide_title") or "",
                }
            )
    return records
