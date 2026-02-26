from __future__ import annotations

import copy
import html
import json
import re
from pathlib import Path
from typing import Any

from . import artwork

try:
    import pptx  # type: ignore
except Exception:
    pptx = None


def _pptx_available() -> bool:
    return pptx is not None


def _clean_text(value: object) -> str:
    return str(value or "").strip()


def _slug_token(value: object) -> str:
    token = re.sub(r"[^a-z0-9]+", "_", _clean_text(value).lower()).strip("_")
    return token or "slide"


def _artwork_ok(payload: dict[str, Any]) -> bool:
    marker = payload.get("ok")
    if isinstance(marker, bool):
        return marker
    token = _clean_text(marker).lower()
    return token in {"true", "1", "yes", "ok"}


def _normalize_snapshot_rel_path(raw_path: object, output_root: Path) -> str:
    token = _clean_text(raw_path).replace("\\", "/")
    if not token:
        return ""
    if token.startswith("./"):
        return token[2:]
    path = Path(token)
    if path.is_absolute():
        try:
            return path.resolve().relative_to(output_root.resolve()).as_posix()
        except Exception:
            return token
    return token


def _materialize_diagram_snapshots(
    slide_ast: dict[str, Any],
    *,
    output_root: Path,
    deck_id: str,
) -> tuple[dict[str, Any], list[dict[str, Any]], list[dict[str, Any]]]:
    prepared = copy.deepcopy(slide_ast or {})
    slides = [item for item in list(prepared.get("slides") or []) if isinstance(item, dict)]
    snapshot_rows: list[dict[str, Any]] = []
    snapshot_errors: list[dict[str, Any]] = []
    deck_token = _slug_token(deck_id)
    for slide_idx, slide in enumerate(slides, start=1):
        slide_id = _clean_text(slide.get("slide_id")) or f"SLIDE-{slide_idx:02d}"
        block_rows = [item for item in list(slide.get("body_blocks") or []) if isinstance(item, dict)]
        for block_idx, block in enumerate(block_rows, start=1):
            if _clean_text(block.get("type")).lower() != "diagram":
                continue
            source = _clean_text(block.get("spec"))
            if not source:
                snapshot_errors.append(
                    {"slide_id": slide_id, "block_index": block_idx, "error": "empty_spec", "message": "diagram spec empty"}
                )
                continue
            engine = _clean_text(block.get("engine")).lower() or "mermaid"
            rel_output = (
                f"deck_assets/{deck_token}/diagrams/"
                f"{_slug_token(slide_id)}_{block_idx:02d}_{_slug_token(engine)}.svg"
            )
            if engine == "d2":
                payload = artwork.render_d2_svg(output_root, source, output_rel_path=rel_output)
            else:
                payload = artwork.render_mermaid_diagram(
                    output_root,
                    source,
                    output_rel_path=rel_output,
                    output_format="svg",
                )
            if _artwork_ok(payload):
                rel_path = _normalize_snapshot_rel_path(payload.get("path"), output_root)
                block["snapshot_path"] = rel_path
                block["snapshot_format"] = "svg"
                block["snapshot_engine"] = engine
                snapshot_rows.append(
                    {
                        "slide_id": slide_id,
                        "block_index": block_idx,
                        "engine": engine,
                        "path": rel_path,
                    }
                )
            else:
                snapshot_errors.append(
                    {
                        "slide_id": slide_id,
                        "block_index": block_idx,
                        "engine": engine,
                        "error": _clean_text(payload.get("error")) or "render_failed",
                        "message": _clean_text(payload.get("message")) or "diagram snapshot render failed",
                    }
                )
    prepared["slides"] = slides
    return prepared, snapshot_rows, snapshot_errors


def _render_block_for_html(block: dict[str, Any]) -> str:
    block_type = _clean_text(block.get("type")).lower()
    if block_type == "text":
        text = html.escape(_clean_text(block.get("text")))
        return f"<p>{text}</p>" if text else ""
    if block_type == "bullets":
        items = [html.escape(str(item).strip()) for item in list(block.get("items") or []) if str(item).strip()]
        if not items:
            return ""
        rows = "\n".join(f"<li>{item}</li>" for item in items)
        return f"<ul>\n{rows}\n</ul>"
    if block_type == "table":
        columns = [html.escape(str(item).strip()) for item in list(block.get("columns") or []) if str(item).strip()]
        rows = [row for row in list(block.get("rows") or []) if isinstance(row, list)]
        head = "".join(f"<th>{item}</th>" for item in columns) if columns else ""
        body_rows: list[str] = []
        for row in rows:
            cols = "".join(f"<td>{html.escape(str(item or ''))}</td>" for item in row)
            body_rows.append(f"<tr>{cols}</tr>")
        body = "\n".join(body_rows)
        if not head and not body:
            return ""
        head_html = f"<thead><tr>{head}</tr></thead>" if head else ""
        body_html = f"<tbody>{body}</tbody>" if body else ""
        return f"<table>{head_html}{body_html}</table>"
    if block_type == "diagram":
        snapshot_path = html.escape(_clean_text(block.get("snapshot_path")))
        engine = html.escape(_clean_text(block.get("engine")) or "diagram")
        spec = html.escape(_clean_text(block.get("spec")))
        if snapshot_path:
            return (
                "<figure>"
                f"<img src=\"{snapshot_path}\" alt=\"diagram snapshot ({engine})\" />"
                f"<figcaption>Diagram snapshot ({engine})</figcaption>"
                "</figure>"
            )
        if not spec:
            return ""
        return f"<pre data-diagram-engine=\"{engine}\">{spec}</pre>"
    if block_type == "chart":
        payload = html.escape(json.dumps(block, ensure_ascii=False))
        return f"<pre data-chart=\"json\">{payload}</pre>"
    if block_type == "image":
        path = html.escape(_clean_text(block.get("path")))
        alt = html.escape(_clean_text(block.get("alt")) or "slide image")
        if path:
            return f"<img src=\"{path}\" alt=\"{alt}\" />"
        return f"<div class=\"image-placeholder\">{alt}</div>"
    payload = html.escape(json.dumps(block, ensure_ascii=False))
    return f"<pre>{payload}</pre>"


def render_slide_ast_html(
    slide_ast: dict[str, Any],
    *,
    output_html: str | Path,
    deck_title: str = "Federlicht Deck",
) -> dict[str, Any]:
    out_path = Path(output_html)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = [item for item in list((slide_ast or {}).get("slides") or []) if isinstance(item, dict)]
    title = _clean_text(deck_title) or "Federlicht Deck"
    lines = [
        "<!doctype html>",
        "<html>",
        "<head>",
        "<meta charset=\"utf-8\" />",
        f"<title>{html.escape(title)}</title>",
        "<style>",
        "body { font-family: Segoe UI, Arial, sans-serif; margin: 24px; background: #f7fafc; color: #111827; }",
        ".slide { background: #ffffff; border: 1px solid #dbeafe; border-radius: 12px; padding: 16px; margin: 0 0 18px; }",
        ".slide h2 { margin: 0 0 6px; }",
        ".subtitle { color: #4b5563; font-size: 13px; margin: 0 0 10px; }",
        "table { border-collapse: collapse; width: 100%; }",
        "th, td { border: 1px solid #cbd5e1; padding: 6px 8px; font-size: 13px; text-align: left; }",
        "pre { white-space: pre-wrap; background: #f8fafc; border: 1px solid #e2e8f0; padding: 10px; border-radius: 8px; }",
        ".image-placeholder { border: 1px dashed #93c5fd; background: #eff6ff; padding: 10px; border-radius: 8px; }",
        "</style>",
        "</head>",
        "<body>",
        f"<h1>{html.escape(title)}</h1>",
    ]
    for entry in slides:
        slide_id = _clean_text(entry.get("slide_id")) or "SLIDE-??"
        title_block = dict(entry.get("title_block") or {})
        headline = _clean_text(title_block.get("headline")) or slide_id
        subheadline = _clean_text(title_block.get("subheadline"))
        lines.append("<section class=\"slide\">")
        lines.append(f"<h2>{html.escape(slide_id)} - {html.escape(headline)}</h2>")
        if subheadline:
            lines.append(f"<div class=\"subtitle\">{html.escape(subheadline)}</div>")
        body_blocks = [item for item in list(entry.get("body_blocks") or []) if isinstance(item, dict)]
        for block in body_blocks:
            rendered = _render_block_for_html(block)
            if rendered:
                lines.append(rendered)
        footer = dict(entry.get("citation_footer") or {})
        refs = [html.escape(str(item).strip()) for item in list(footer.get("refs") or []) if str(item).strip()]
        if refs:
            lines.append("<div class=\"subtitle\">Refs: " + ", ".join(refs[:5]) + "</div>")
        lines.append("</section>")
    lines.extend(["</body>", "</html>"])
    out_path.write_text("\n".join(lines), encoding="utf-8")
    return {
        "ok": True,
        "path": out_path.as_posix(),
        "slide_count": len(slides),
        "bytes": out_path.stat().st_size if out_path.exists() else 0,
    }


def _extract_pptx_text_lines(slide_entry: dict[str, Any]) -> list[str]:
    lines: list[str] = []
    body_blocks = [item for item in list(slide_entry.get("body_blocks") or []) if isinstance(item, dict)]
    for block in body_blocks:
        block_type = _clean_text(block.get("type")).lower()
        if block_type == "text":
            text = _clean_text(block.get("text"))
            if text:
                lines.append(text)
            continue
        if block_type == "bullets":
            for item in list(block.get("items") or []):
                token = _clean_text(item)
                if token:
                    lines.append(f"- {token}")
            continue
        if block_type == "table":
            columns = [str(item).strip() for item in list(block.get("columns") or []) if str(item).strip()]
            if columns:
                lines.append(" | ".join(columns))
            for row in list(block.get("rows") or []):
                if isinstance(row, list):
                    lines.append(" | ".join(str(item) for item in row))
            continue
        if block_type == "diagram":
            engine = _clean_text(block.get("engine") or "mermaid")
            snapshot_path = _clean_text(block.get("snapshot_path"))
            if snapshot_path:
                lines.append(f"[diagram:{engine}] snapshot={snapshot_path}")
            else:
                lines.append(f"[diagram] {engine}")
            continue
        if block_type == "chart":
            lines.append("[chart] data attached")
            continue
        if block_type == "image":
            alt = _clean_text(block.get("alt")) or "image placeholder"
            lines.append(f"[image] {alt}")
            continue
    return lines


def render_slide_ast_pptx(
    slide_ast: dict[str, Any],
    *,
    output_pptx: str | Path,
) -> dict[str, Any]:
    out_path = Path(output_pptx)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = [item for item in list((slide_ast or {}).get("slides") or []) if isinstance(item, dict)]
    if not _pptx_available():
        return {
            "ok": False,
            "reason": "python-pptx not installed",
            "path": out_path.as_posix(),
            "slide_count": len(slides),
            "bytes": 0,
        }
    presentation = pptx.Presentation()  # type: ignore[attr-defined]
    layout_title_and_content = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
    for entry in slides:
        slide_obj = presentation.slides.add_slide(layout_title_and_content)
        title_block = dict(entry.get("title_block") or {})
        headline = _clean_text(title_block.get("headline")) or _clean_text(entry.get("slide_id")) or "Untitled Slide"
        if getattr(slide_obj.shapes, "title", None) is not None:
            slide_obj.shapes.title.text = headline
        content_placeholder = None
        if len(slide_obj.placeholders) > 1:
            content_placeholder = slide_obj.placeholders[1]
        lines = _extract_pptx_text_lines(entry)
        text_payload = "\n".join(lines).strip()
        if content_placeholder is not None and text_payload:
            try:
                content_placeholder.text = text_payload
            except Exception:
                pass
    presentation.save(out_path)
    return {
        "ok": True,
        "path": out_path.as_posix(),
        "slide_count": len(slides),
        "bytes": out_path.stat().st_size if out_path.exists() else 0,
    }


def render_slide_ast_bundle(
    slide_ast: dict[str, Any],
    *,
    output_dir: str | Path,
    deck_id: str,
    deck_title: str = "",
    export_html: bool = True,
    export_pptx: bool = True,
) -> dict[str, Any]:
    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    token = _clean_text(deck_id) or "deck"
    html_path = out_dir / f"{token}.html"
    pptx_path = out_dir / f"{token}.pptx"
    prepared_ast, snapshot_rows, snapshot_errors = _materialize_diagram_snapshots(
        slide_ast,
        output_root=out_dir,
        deck_id=token,
    )
    results: dict[str, Any] = {
        "deck_id": token,
        "output_dir": out_dir.as_posix(),
        "diagram_snapshot_count": len(snapshot_rows),
        "diagram_snapshot_paths": [str(item.get("path") or "") for item in snapshot_rows if _clean_text(item.get("path"))],
        "diagram_snapshot_errors": snapshot_errors,
        "html": {"ok": False, "path": html_path.as_posix(), "slide_count": 0, "bytes": 0},
        "pptx": {"ok": False, "path": pptx_path.as_posix(), "slide_count": 0, "bytes": 0},
    }
    if export_html:
        results["html"] = render_slide_ast_html(
            prepared_ast,
            output_html=html_path,
            deck_title=deck_title or token,
        )
    if export_pptx:
        results["pptx"] = render_slide_ast_pptx(prepared_ast, output_pptx=pptx_path)
    return results
